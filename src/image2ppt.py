#!/usr/bin/env python3
import sys
import argparse
import shutil
from pathlib import Path
import collections.abc  # noqa: F401 (to avoid pptx import error)

from pptx import Presentation
from pptx.dml.color import RGBColor


def add_image(prs, image):
    """Add `image` to `prs` as a single slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = RGBColor(0, 0, 0)
    slide.shapes.add_picture(image, 0, 0,
                             width=prs.slide_width,
                             height=prs.slide_height)


def append_pptx(in_pptx, out_pptx, images):
    """Append `images` to `in_pptx`, and save the result to `out_pptx`."""
    in_pptx = Path(in_pptx)
    out_pptx = Path(out_pptx)
    if out_pptx.suffix != '.pptx':
        raise ValueError(f'Output file "{out_pptx}" is not a .pptx file')
    if not out_pptx.exists() or not out_pptx.samefile(in_pptx):
        shutil.copyfile(in_pptx, out_pptx)
    prs = Presentation(out_pptx)
    for image in images:
        add_image(prs, image)
    prs.save(out_pptx)


def main(args):
    parser = argparse.ArgumentParser()
    parser.add_argument('in_pptx', help='Input PPTX file')
    parser.add_argument('out_pptx', help='Output PPTX file')
    parser.add_argument('images', nargs='+',
                        help='List of images to be appended to the PPTX file')
    args = parser.parse_args(args)
    append_pptx(args.in_pptx, args.out_pptx, tuple(args.images))


if __name__ == '__main__':
    main(sys.argv[1:])
